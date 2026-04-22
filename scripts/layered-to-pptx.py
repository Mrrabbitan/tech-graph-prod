#!/usr/bin/env python3
"""Generate an editable PPTX of a multi-layer technical architecture diagram.

Reuses the SAME JSON schema as `layered-to-drawio.py` so one config can produce
both an editable drawio XML and an editable PowerPoint file with identical visual
style (colors, sizing, layout).

Every element on the slide is a native PowerPoint shape (rounded rects, connectors,
text boxes) - no embedded raster image. Open in PowerPoint / WPS / Keynote to edit.

Dependencies:
    pip install python-pptx lxml

Usage:
    python3 scripts/layered-to-pptx.py \\
        --config fixtures/layered-architecture-example.json \\
        --output output/my-arch.pptx
"""
from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Emu, Pt, Inches
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.oxml.ns import qn
    from lxml import etree
except ImportError as e:
    raise SystemExit(
        f"Missing dependency: {e.name}. Install with:\n"
        f"    pip install python-pptx lxml"
    )


def rgb(hex_str: str) -> RGBColor:
    hex_str = hex_str.lstrip('#')
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


def _chip_width_px(text: str) -> int:
    w = 22
    for ch in text:
        if ord(ch) > 0x2E80:
            w += 14
        elif ch.isupper():
            w += 8
        elif ch.isdigit():
            w += 7
        else:
            w += 6.5
    return max(int(w), 58)


def build_pptx(cfg: dict, out_path: Path) -> None:
    page = cfg.get('page', {})
    VCANVAS_W = float(page.get('width', 1440))
    VCANVAS_H = float(page.get('height', 1050))

    title_str = cfg.get('title', 'Technical Architecture')
    subtitle_str = cfg.get('subtitle', '')
    footer_str = cfg.get('footer', '')
    layers = cfg['layers']
    arrows_list = cfg.get('arrows', [])

    SLIDE_W_EMU = int(Inches(13.333))
    SLIDE_H_EMU = int(Inches(7.5))

    def vx(x):
        return int(x / VCANVAS_W * SLIDE_W_EMU)

    def vy(y):
        return int(y / VCANVAS_H * SLIDE_H_EMU)

    def vw(w):
        return int(w / VCANVAS_W * SLIDE_W_EMU)

    def vh(h):
        return int(h / VCANVAS_H * SLIDE_H_EMU)

    MARGIN_X = 30
    CONTENT_W = int(VCANVAS_W) - 2 * MARGIN_X
    LAYER_H = int(cfg.get('layer_height', 168))
    LAYER_GAP = int(cfg.get('layer_gap', 14))
    LAYER_START_Y = 86
    LEFT_STRIP_W = 62
    LABEL_X = MARGIN_X + 10 + LEFT_STRIP_W + 10
    LABEL_W = 258
    COMP_START_X = LABEL_X + LABEL_W + 20
    COMP_AREA_W = (MARGIN_X + CONTENT_W) - COMP_START_X
    COMP_GAP = 14
    COMP_H = 62
    CHIP_H = 22
    CHIP_Y_OFFSET = 132

    prs = Presentation()
    prs.slide_width = SLIDE_W_EMU
    prs.slide_height = SLIDE_H_EMU
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W_EMU, SLIDE_H_EMU)
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb('#ffffff')
    bg.line.fill.background()
    bg.shadow.inherit = False

    def set_text(tf, segments, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
        tf.word_wrap = True
        tf.margin_left = Emu(0)
        tf.margin_right = Emu(0)
        tf.margin_top = Emu(0)
        tf.margin_bottom = Emu(0)
        tf.vertical_anchor = anchor
        for i, seg in enumerate(segments):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            if i == 0:
                for r in list(p.runs):
                    r._r.getparent().remove(r._r)
                p.text = ''
            p.alignment = align
            runs = seg if isinstance(seg, list) else [seg]
            for text, props in runs:
                run = p.add_run()
                run.text = text
                f = run.font
                if 'size' in props:
                    f.size = Pt(props['size'])
                if props.get('bold'):
                    f.bold = True
                if props.get('italic'):
                    f.italic = True
                if 'color' in props:
                    f.color.rgb = rgb(props['color'])
                f.name = props.get('name', 'PingFang SC')

    def rounded_rect(x, y, w, h, fill, border=None, radius=0.12):
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, vx(x), vy(y), vw(w), vh(h))
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
        shp.fill.solid()
        shp.fill.fore_color.rgb = rgb(fill)
        if border:
            shp.line.color.rgb = rgb(border)
            shp.line.width = Pt(1.25)
        else:
            shp.line.fill.background()
        shp.shadow.inherit = False
        return shp

    def plain_rect(x, y, w, h, fill):
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, vx(x), vy(y), vw(w), vh(h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = rgb(fill)
        shp.line.fill.background()
        shp.shadow.inherit = False
        return shp

    def text_box(x, y, w, h, segments, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
        tb = slide.shapes.add_textbox(vx(x), vy(y), vw(w), vh(h))
        set_text(tb.text_frame, segments, align=align, anchor=anchor)
        return tb

    # Title
    text_box(
        0, 18, VCANVAS_W, 32,
        [[(title_str, {'size': 22, 'bold': True, 'color': '#0f172a'})]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    if subtitle_str:
        text_box(
            0, 50, VCANVAS_W, 22,
            [[(subtitle_str, {'size': 13, 'color': '#64748b'})]],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )

    layer_top_ys = []
    for i, layer in enumerate(layers):
        top = LAYER_START_Y + i * (LAYER_H + LAYER_GAP)
        bottom = top + LAYER_H
        layer_top_ys.append((top, bottom))

        bg_color = layer.get('bg', '#f8fafc')
        border_color = layer.get('border', '#cbd5e1')
        accent = layer.get('accent', '#334155')
        idx_str = layer.get('idx', str(i + 1).zfill(2))

        rounded_rect(MARGIN_X, top, CONTENT_W, LAYER_H, bg_color, border=border_color, radius=0.04)
        strip = rounded_rect(MARGIN_X + 10, top + 14, LEFT_STRIP_W, LAYER_H - 28, accent, radius=0.1)
        set_text(
            strip.text_frame,
            [[(idx_str, {'size': 26, 'bold': True, 'color': '#ffffff'})]],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )

        text_box(
            LABEL_X, top + 12, LABEL_W, 28,
            [[(layer.get('name', ''), {'size': 17, 'bold': True, 'color': accent})]],
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
        )
        if layer.get('en'):
            text_box(
                LABEL_X, top + 40, LABEL_W, 18,
                [[(layer['en'], {'size': 10.5, 'italic': True, 'color': '#64748b', 'name': 'Helvetica'})]],
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
            )

        hl = rounded_rect(LABEL_X, top + 62, LABEL_W, 90, '#ffffff', border=border_color, radius=0.08)
        hl.text_frame.margin_left = Emu(int(Pt(6)))
        hl.text_frame.margin_right = Emu(int(Pt(6)))
        hl.text_frame.margin_top = Emu(int(Pt(5)))
        hl.text_frame.margin_bottom = Emu(int(Pt(5)))
        set_text(
            hl.text_frame,
            [
                [('技术亮点', {'size': 10, 'bold': True, 'color': '#64748b'})],
                [(layer.get('highlights', ''), {'size': 11, 'color': '#1e293b'})],
            ],
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        )

        components = layer.get('components', [])
        num = max(len(components), 1)
        comp_w = (COMP_AREA_W - (num - 1) * COMP_GAP) / num
        for j, comp in enumerate(components):
            cx = COMP_START_X + j * (comp_w + COMP_GAP)
            cy = top + 20
            shp = rounded_rect(cx, cy, comp_w, COMP_H, '#ffffff', border=accent, radius=0.15)
            shp.line.width = Pt(1.4)
            shp.text_frame.margin_left = Emu(int(Pt(4)))
            shp.text_frame.margin_right = Emu(int(Pt(4)))
            shp.text_frame.margin_top = Emu(int(Pt(4)))
            shp.text_frame.margin_bottom = Emu(int(Pt(4)))
            set_text(
                shp.text_frame,
                [
                    [(comp.get('title', ''), {'size': 12, 'bold': True, 'color': accent})],
                    [(comp.get('subtitle', ''), {'size': 9.5, 'color': '#475569'})],
                ],
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
            )

        text_box(
            COMP_START_X, top + CHIP_Y_OFFSET - 2, 60, 20,
            [[('技术栈', {'size': 10, 'bold': True, 'color': '#64748b'})]],
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
        )
        chip_x = COMP_START_X + 56
        chip_y = top + CHIP_Y_OFFSET
        for chip_text in layer.get('chips', []):
            w = _chip_width_px(chip_text)
            chip = rounded_rect(chip_x, chip_y, w, CHIP_H, accent, radius=0.5)
            chip.text_frame.margin_left = Emu(int(Pt(4)))
            chip.text_frame.margin_right = Emu(int(Pt(4)))
            chip.text_frame.margin_top = Emu(0)
            chip.text_frame.margin_bottom = Emu(0)
            set_text(
                chip.text_frame,
                [[(chip_text, {'size': 10, 'bold': True, 'color': '#ffffff'})]],
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
            )
            chip_x += w + 8

    # Inter-layer bidirectional connectors
    arrow_x_center = VCANVAS_W / 2
    for idx in range(len(layers) - 1):
        src_bottom = layer_top_ys[idx][1]
        tgt_top = layer_top_ys[idx + 1][0]
        connector = slide.shapes.add_connector(
            1, vx(arrow_x_center), vy(src_bottom), vx(arrow_x_center), vy(tgt_top),
        )
        line = connector.line
        line.color.rgb = rgb('#475569')
        line.width = Pt(2)
        ln = connector.line._get_or_add_ln()
        head_end = etree.SubElement(ln, qn('a:headEnd'))
        head_end.set('type', 'triangle')
        head_end.set('w', 'med')
        head_end.set('len', 'med')
        tail_end = etree.SubElement(ln, qn('a:tailEnd'))
        tail_end.set('type', 'triangle')
        tail_end.set('w', 'med')
        tail_end.set('len', 'med')

        label_text = arrows_list[idx] if idx < len(arrows_list) else ''
        if label_text:
            mid_y = (src_bottom + tgt_top) / 2
            lbl_w = 170
            lbl_h = 18
            lbl_x = arrow_x_center - lbl_w / 2
            bg_box = plain_rect(lbl_x, mid_y - lbl_h / 2, lbl_w, lbl_h, '#ffffff')
            set_text(
                bg_box.text_frame,
                [[(label_text, {'size': 10, 'color': '#334155'})]],
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
            )

    if footer_str:
        footer_y = VCANVAS_H - 44
        shp = rounded_rect(MARGIN_X, footer_y, CONTENT_W, 30, '#f1f5f9', border='#cbd5e1', radius=0.1)
        shp.text_frame.margin_left = Emu(int(Pt(8)))
        shp.text_frame.margin_right = Emu(int(Pt(8)))
        shp.text_frame.margin_top = Emu(0)
        shp.text_frame.margin_bottom = Emu(0)
        set_text(
            shp.text_frame,
            [[(footer_str, {'size': 11, 'color': '#475569'})]],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )

    # Speaker notes
    notes_lines = [f'{layer.get("idx", str(i+1).zfill(2))} · {layer.get("name", "")} — ' +
                   ' / '.join(layer.get('chips', []))
                   for i, layer in enumerate(layers)]
    if footer_str:
        notes_lines.append('')
        notes_lines.append(footer_str)
    slide.notes_slide.notes_text_frame.text = '\n'.join(notes_lines)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))


def load_config(args) -> dict:
    if args.config:
        with open(args.config, 'r', encoding='utf-8') as f:
            return json.load(f)
    data = sys.stdin.read().strip()
    if not data:
        raise SystemExit('No JSON config provided (use --config or stdin).')
    return json.loads(data)


def main() -> int:
    ap = argparse.ArgumentParser(description='Generate layered-architecture editable PPTX.')
    ap.add_argument('--config', '-c', help='Path to JSON config file.')
    ap.add_argument('--output', '-o', required=True, help='Output .pptx path.')
    args = ap.parse_args()

    cfg = load_config(args)
    out = Path(args.output)
    build_pptx(cfg, out)
    print(f'✓ pptx written: {out}')
    return 0


if __name__ == '__main__':
    raise SystemExit(main())
